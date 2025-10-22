#!/usr/bin/env python3
import argparse
from pathlib import Path
from PIL import Image
from pptx import Presentation
from pptx.util import Inches
import re

def load_abstract(text_path: Path, max_bullets=5):
    if not text_path.exists():
        return []
    text = text_path.read_text(encoding="utf-8", errors="ignore")
    m = re.search(r"(?im)^(Abstract|Resumen)\s*$", text)
    if m:
        start = m.end()
        m2 = re.search(r"(?im)^(Introduction|Introducción|Background|Methods|Resultados|Results)\s*$", text[start:])
        end = start + (m2.start() if m2 else len(text))
        segment = text[start:end]
    else:
        segment = text[:1200]
    parts = re.split(r"(?<=[.!?])\s+", segment.strip())
    parts = [p.strip() for p in parts if p.strip()]
    bullets = []
    for p in parts:
        if len(p) > 300:
            p = p[:280].rstrip() + "..."
        bullets.append(p)
        if len(bullets) >= max_bullets:
            break
    return bullets

def add_title_slide(prs, title, subtitle=""):
    layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(layout)
    slide.shapes.title.text = title
    slide.placeholders[1].text = subtitle

def add_bullets_slide(prs, title, bullets):
    layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(layout)
    slide.shapes.title.text = title
    tf = slide.placeholders[1].text_frame
    tf.clear()
    for i, b in enumerate(bullets):
        if i == 0:
            tf.text = b
        else:
            p = tf.add_paragraph()
            p.text = b
            p.level = 0

def add_image_slide(prs, image_path: Path, margin_in=0.3):
    blank = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank)
    slide_w = prs.slide_width
    slide_h = prs.slide_height

    slide_w_in = slide_w / 914400
    slide_h_in = slide_h / 914400

    with Image.open(image_path) as im:
        img_w_px, img_h_px = im.size
        dpi = im.info.get("dpi", (96, 96))[0] or 96
        img_w_in = img_w_px / dpi
        img_h_in = img_h_px / dpi

    avail_w = max(slide_w_in - 2*margin_in, 0.1)
    avail_h = max(slide_h_in - 2*margin_in, 0.1)

    scale = min(avail_w / img_w_in, avail_h / img_h_in)
    target_w_in = img_w_in * scale
    target_h_in = img_h_in * scale

    left_in = (slide_w_in - target_w_in) / 2
    top_in = (slide_h_in - target_h_in) / 2

    slide.shapes.add_picture(
        str(image_path),
        Inches(left_in),
        Inches(top_in),
        width=Inches(target_w_in),
        height=Inches(target_h_in),
    )

def extract_num(path: Path):
    m = re.search(r"(\d+)", path.stem)
    return int(m.group(1)) if m else None

def main():
    ap = argparse.ArgumentParser(description="Construye un PPTX a partir de imágenes (una por diapositiva).")
    ap.add_argument("--images-dir", default="assets/pages", help="Directorio con PNG/JPG.")
    ap.add_argument("--glob", default="*.png", help="Patrón de archivo, ej. page-*.png")
    ap.add_argument("--out", default="slides/paper_pages.pptx", help="Ruta de salida PPTX")
    ap.add_argument("--title", default="Presentación del paper", help="Título para la primera diapositiva")
    ap.add_argument("--subtitle", default="", help="Subtítulo para la primera diapositiva")
    ap.add_argument("--abstract", default="analysis/paper.txt", help="Ruta al texto para 'Resumen' (opcional)")
    ap.add_argument("--no-title", action="store_true", help="No agregar diapositiva de título")
    ap.add_argument("--no-abstract", action="store_true", help="No agregar diapositiva de Resumen")
    ap.add_argument("--margin", type=float, default=0.3, help="Margen alrededor de la imagen (pulgadas)")
    ap.add_argument("--start", type=int, default=None, help="Número inicial (inclusive) para filtrar imágenes por número en el nombre")
    ap.add_argument("--end", type=int, default=None, help="Número final (inclusive) para filtrar imágenes por número en el nombre")
    args = ap.parse_args()

    img_dir = Path(args.images_dir)
    imgs = sorted(img_dir.glob(args.glob), key=lambda p: (extract_num(p) if extract_num(p) is not None else float("inf"), p.name))

    if args.start is not None or args.end is not None:
        s = args.start if args.start is not None else -float("inf")
        e = args.end if args.end is not None else float("inf")
        imgs = [p for p in imgs if (extract_num(p) is not None and s <= extract_num(p) <= e)]

    if not imgs:
        raise SystemExit(f"No se encontraron imágenes con patrón {args.glob} en {img_dir} con el rango solicitado.")

    prs = Presentation()

    if not args.no_title:
        add_title_slide(prs, args.title, args.subtitle)

    if not args.no_abstract and args.abstract:
        bullets = load_abstract(Path(args.abstract))
        if bullets:
            add_bullets_slide(prs, "Resumen (digest)", bullets)

    for img in imgs:
        add_image_slide(prs, img, margin_in=args.margin)

    out_path = Path(args.out)
    out_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(out_path))
    print(f"OK: generado {out_path} con {len(prs.slides)} diapositivas.")

if __name__ == "__main__":
    main()
